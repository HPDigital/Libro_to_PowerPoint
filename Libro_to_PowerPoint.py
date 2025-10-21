"""
Libro_to_PowerPoint
"""

#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import openai
from pptx import Presentation

# Configura tu clave de API de OpenAI
openai.api_key = 'YOUR_API_KEY_HERE'

# Función para obtener el resumen de un capítulo
def obtener_resumen(texto_capitulo):
    prompt_resumen = f"Por favor, genera un resumen detallado del siguiente capítulo del libro: {texto_capitulo}"

    response = openai.Completion.create(
        engine="text-davinci-003",  # O el modelo que prefieras usar
        prompt=prompt_resumen,
        max_tokens=1500
    )

    return response['choices'][0]['text']

# Función para generar la estructura de las diapositivas
def generar_estructura_diapositivas(resumen):
    prompt_estructura = f"Con base en el siguiente resumen, crea una estructura para una presentación de 10 diapositivas: {resumen}"

    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt_estructura,
        max_tokens=1000
    )

    return response['choices'][0]['text']

# Función para generar el contenido de cada diapositiva
def generar_contenido_diapositiva(punto_clave):
    prompt_contenido = f"Genera el contenido de una diapositiva en formato claro y conciso para este punto clave: {punto_clave}"

    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt_contenido,
        max_tokens=300
    )

    return response['choices'][0]['text']

# Función para crear la presentación en PowerPoint
def crear_presentacion(capitulo, estructura):
    prs = Presentation()
    for i, punto in enumerate(estructura):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Layout de título y contenido
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]

        title.text = f"Diapositiva {i+1}: {punto['titulo']}"
        body.text = punto['contenido']

    prs.save(f"Presentacion_Capitulo_{capitulo}.pptx")

# Función principal para procesar cada capítulo
def procesar_capitulo(capitulo_num, texto_capitulo):
    print(f"Procesando capítulo {capitulo_num}...")

    # Obtener resumen del capítulo
    resumen = obtener_resumen(texto_capitulo)
    print(f"Resumen del capítulo {capitulo_num}: {resumen}")

    # Generar la estructura de las diapositivas
    estructura = generar_estructura_diapositivas(resumen)
    puntos_claves = estructura.split("\n")

    # Generar contenido para cada diapositiva
    diapositivas = []
    for punto in puntos_claves:
        contenido = generar_contenido_diapositiva(punto)
        diapositivas.append({"titulo": punto, "contenido": contenido})

    # Crear la presentación
    crear_presentacion(capitulo_num, diapositivas)

# Iterar sobre cada capítulo
def procesar_libro(lista_capitulos):
    for i, capitulo in enumerate(lista_capitulos):
        procesar_capitulo(i + 1, capitulo)

# Aquí deberías cargar el libro y separar los capítulos en una lista
# En este ejemplo, se asume que los capítulos ya están en una lista:
lista_capitulos = [
    "Texto del capítulo 1...",
    "Texto del capítulo 2...",
    # Añadir más capítulos según sea necesario
]

# Ejecutar el procesamiento del libro
procesar_libro(lista_capitulos)



if __name__ == "__main__":
    pass
